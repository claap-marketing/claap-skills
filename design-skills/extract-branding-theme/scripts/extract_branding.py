#!/usr/bin/env python3
"""
extract_branding.py — Reverse-engineer a .pptx into a structured JSON design system.

Usage:
    python extract_branding.py <path_to_pptx> [--output-dir /mnt/user-data/outputs]

Outputs:
    - branding_template.json  (full design system)
    - assets/                 (extracted images with classification)
"""

import sys
import os
import json
import zipfile
import hashlib
import collections
import re
from pathlib import Path
from copy import deepcopy

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
    from pptx.oxml.ns import qn
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx --break-system-packages")
    sys.exit(1)

try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
DPI = 96
EMU_PER_INCH = 914400
EMU_PER_PT = 12700

NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}


def emu_to_inches(emu):
    if emu is None:
        return None
    return round(emu / EMU_PER_INCH, 4)


def emu_to_px(emu):
    if emu is None:
        return None
    return round(emu / EMU_PER_INCH * DPI, 1)


def emu_to_pt(emu):
    if emu is None:
        return None
    return round(emu / EMU_PER_PT, 1)


def dim(emu):
    """Return a dimension dict with inches, px, emu."""
    if emu is None:
        return {"inches": None, "px": None, "emu": None}
    return {
        "inches": emu_to_inches(emu),
        "px": emu_to_px(emu),
        "emu": int(emu),
    }


def rgb_to_hex(rgb):
    """Convert pptx RGBColor or string to #RRGGBB."""
    if rgb is None:
        return None
    s = str(rgb).upper()
    if s.startswith('#'):
        return s
    if len(s) == 6 and all(c in '0123456789ABCDEF' for c in s):
        return f"#{s}"
    return None


# ---------------------------------------------------------------------------
# Phase 1: Theme Colors
# ---------------------------------------------------------------------------
def extract_theme_colors(prs):
    """Extract theme color scheme from the first slide master."""
    theme_colors = {}
    try:
        master_el = prs.slide_masters[0].element
        clr_scheme = master_el.find('.//' + qn('a:clrScheme'))
        if clr_scheme is None:
            return theme_colors
        for child in clr_scheme:
            tag = child.tag.split('}')[-1]
            srgb = child.find(qn('a:srgbClr'))
            sys_clr = child.find(qn('a:sysClr'))
            if srgb is not None:
                theme_colors[tag] = f"#{srgb.get('val', '000000').upper()}"
            elif sys_clr is not None:
                val = sys_clr.get('lastClr', sys_clr.get('val', '000000'))
                theme_colors[tag] = f"#{val.upper()}"
    except Exception:
        pass
    return theme_colors


def extract_theme_fonts(prs):
    """Extract major/minor theme font names."""
    major = None
    minor = None
    try:
        master_el = prs.slide_masters[0].element
        # major font
        mj = master_el.find('.//' + qn('a:majorFont'))
        if mj is not None:
            latin = mj.find(qn('a:latin'))
            if latin is not None:
                major = latin.get('typeface')
        # minor font
        mn = master_el.find('.//' + qn('a:minorFont'))
        if mn is not None:
            latin = mn.find(qn('a:latin'))
            if latin is not None:
                minor = latin.get('typeface')
    except Exception:
        pass
    return major, minor


def resolve_font_name(name, major_font, minor_font):
    """Resolve theme font references like +mj-lt, +mn-lt."""
    if name is None:
        return None
    if name.startswith('+mj') or name == 'Heading':
        return major_font or name
    if name.startswith('+mn') or name == 'Body':
        return minor_font or name
    return name


def safe_color_from_run(run):
    """Safely extract hex color from a text run."""
    try:
        if run.font.color and run.font.color.rgb:
            return rgb_to_hex(run.font.color.rgb)
    except Exception:
        pass
    return None


def safe_fill_color(shape):
    """Safely extract fill color from a shape."""
    try:
        fill = shape.fill
        if fill.type is not None:
            if hasattr(fill, 'fore_color') and fill.fore_color and fill.fore_color.rgb:
                return rgb_to_hex(fill.fore_color.rgb)
    except Exception:
        pass
    # Try XML fallback
    try:
        solid = shape.element.find('.//' + qn('a:solidFill'))
        if solid is not None:
            srgb = solid.find(qn('a:srgbClr'))
            if srgb is not None:
                return f"#{srgb.get('val', '').upper()}"
    except Exception:
        pass
    return None


def safe_line_color(shape):
    """Safely extract line/border color from a shape."""
    try:
        ln = shape.element.find(qn('p:spPr') + '/' + qn('a:ln'))
        if ln is None:
            ln = shape.element.find('.//' + qn('a:ln'))
        if ln is not None:
            solid = ln.find(qn('a:solidFill'))
            if solid is not None:
                srgb = solid.find(qn('a:srgbClr'))
                if srgb is not None:
                    return f"#{srgb.get('val', '').upper()}"
    except Exception:
        pass
    return None


def collect_all_colors(prs):
    """Walk every shape and collect fill, font, and border colors with frequency."""
    fill_counter = collections.Counter()
    font_counter = collections.Counter()
    border_counter = collections.Counter()

    for slide in prs.slides:
        for shape in slide.shapes:
            # Fill
            fc = safe_fill_color(shape)
            if fc:
                fill_counter[fc] += 1
            # Border
            bc = safe_line_color(shape)
            if bc:
                border_counter[bc] += 1
            # Text colors
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        tc = safe_color_from_run(run)
                        if tc:
                            font_counter[tc] += 1
            # Table cells
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                tc = safe_color_from_run(run)
                                if tc:
                                    font_counter[tc] += 1

    return fill_counter, font_counter, border_counter


def extract_gradients(prs):
    """Extract gradient fills from shapes (XML-level)."""
    gradients = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            try:
                grad = shape.element.find('.//' + qn('a:gradFill'))
                if grad is None:
                    continue
                gs_list = grad.findall(qn('a:gsLst') + '/' + qn('a:gs'))
                stops = []
                for gs in gs_list:
                    pos = int(gs.get('pos', '0')) / 100000.0
                    srgb = gs.find('.//' + qn('a:srgbClr'))
                    hex_val = f"#{srgb.get('val', '').upper()}" if srgb is not None else None
                    stops.append({"hex": hex_val, "position": round(pos, 3)})
                lin = grad.find(qn('a:lin'))
                angle = int(lin.get('ang', '0')) / 60000 if lin is not None else None
                gradients.append({
                    "type": "linear" if lin is not None else "radial",
                    "angle": angle,
                    "stops": stops,
                    "used_on": f"slide {slide_idx} shape"
                })
            except Exception:
                continue
    return gradients


def build_dominant_colors(fill_counter, font_counter, border_counter):
    """Merge all color counters and rank by frequency."""
    merged = collections.Counter()
    usage_map = collections.defaultdict(set)
    for c, n in fill_counter.items():
        merged[c] += n
        usage_map[c].add("fill")
    for c, n in font_counter.items():
        merged[c] += n
        usage_map[c].add("text")
    for c, n in border_counter.items():
        merged[c] += n
        usage_map[c].add("border")

    dominant = []
    for color, freq in merged.most_common(30):
        dominant.append({
            "hex": color,
            "frequency": freq,
            "usage": " | ".join(sorted(usage_map[color])),
        })
    return dominant


# ---------------------------------------------------------------------------
# Phase 2: Typography
# ---------------------------------------------------------------------------
def extract_typography(prs, major_font, minor_font):
    """Extract and cluster all text styles."""
    style_counter = collections.Counter()
    style_examples = {}

    for slide in prs.slides:
        for shape in slide.shapes:
            frames = []
            if shape.has_text_frame:
                frames.append(shape.text_frame)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        frames.append(cell.text_frame)

            for tf in frames:
                for para in tf.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        font_name = resolve_font_name(run.font.name, major_font, minor_font)
                        font_size = round(run.font.size / EMU_PER_PT, 1) if run.font.size else None
                        bold = run.font.bold
                        italic = run.font.italic
                        underline = run.font.underline
                        color = safe_color_from_run(run)

                        alignment = None
                        try:
                            alignment = str(para.alignment).split('(')[0] if para.alignment else None
                        except Exception:
                            pass

                        line_spacing = None
                        space_before = None
                        space_after = None
                        try:
                            pf = para._pPr
                            if pf is not None:
                                lnSpc = pf.find(qn('a:lnSpc'))
                                if lnSpc is not None:
                                    spc_pts = lnSpc.find('.//' + qn('a:spcPts'))
                                    if spc_pts is not None:
                                        line_spacing = int(spc_pts.get('val', '0')) / 100
                                spcBef = pf.find(qn('a:spcBef'))
                                if spcBef is not None:
                                    pts = spcBef.find('.//' + qn('a:spcPts'))
                                    if pts is not None:
                                        space_before = int(pts.get('val', '0')) / 100
                                spcAft = pf.find(qn('a:spcAft'))
                                if spcAft is not None:
                                    pts = spcAft.find('.//' + qn('a:spcPts'))
                                    if pts is not None:
                                        space_after = int(pts.get('val', '0')) / 100
                        except Exception:
                            pass

                        key = (font_name, font_size, bold, italic, underline, color, alignment)
                        style_counter[key] += 1
                        if key not in style_examples:
                            style_examples[key] = run.text.strip()[:80]

    # Cluster and assign roles
    styles_ranked = style_counter.most_common()
    text_styles = []
    size_groups = collections.defaultdict(list)

    for (font_name, font_size, bold, italic, underline, color, alignment), freq in styles_ranked:
        entry = {
            "role": "unknown",
            "font_name": font_name,
            "font_size_pt": font_size,
            "bold": bold if bold is not None else False,
            "italic": italic if italic is not None else False,
            "underline": underline if underline is not None else False,
            "color": color,
            "alignment": alignment,
            "frequency": freq,
            "example_text": style_examples.get(
                (font_name, font_size, bold, italic, underline, color, alignment), ""
            ),
        }
        text_styles.append(entry)
        if font_size:
            size_groups[font_size].append(entry)

    # Auto-label by size ranking
    if text_styles:
        unique_sizes = sorted(set(s["font_size_pt"] for s in text_styles if s["font_size_pt"]), reverse=True)
        size_role_map = {}
        role_names = ["title", "heading", "subheading", "body", "caption", "footnote"]
        for i, sz in enumerate(unique_sizes):
            if i < len(role_names):
                size_role_map[sz] = role_names[i]
            else:
                size_role_map[sz] = f"style_{sz}pt"

        # Override: the most frequent style is likely body
        if styles_ranked:
            most_freq_key = styles_ranked[0][0]
            most_freq_size = most_freq_key[1]
            if most_freq_size and most_freq_size in size_role_map:
                size_role_map[most_freq_size] = "body"

        for s in text_styles:
            if s["font_size_pt"] in size_role_map:
                s["role"] = size_role_map[s["font_size_pt"]]

        # Bold variants
        for s in text_styles:
            if s["bold"] and s["role"] == "body":
                s["role"] = "body_bold"

    return text_styles


# ---------------------------------------------------------------------------
# Phase 3: Layout & Spatial
# ---------------------------------------------------------------------------
def extract_layouts(prs):
    """Analyze shape positions to derive margins, grid, and zones."""
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    lefts = []
    tops = []
    rights = []
    bottoms = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.left is not None and shape.width is not None:
                lefts.append(shape.left)
                tops.append(shape.top)
                rights.append(slide_w - shape.left - shape.width)
                bottoms.append(slide_h - shape.top - shape.height)

    content_margins = {
        "left": dim(min(lefts) if lefts else 0),
        "top": dim(min(tops) if tops else 0),
        "right": dim(min(rights) if rights else 0),
        "bottom": dim(min(bottoms) if bottoms else 0),
    }

    # Simple gutter detection: find gaps between horizontally adjacent shapes per slide
    gutters = []
    for slide in prs.slides:
        shapes_sorted = sorted(
            [s for s in slide.shapes if s.left is not None and s.width is not None],
            key=lambda s: s.left
        )
        for i in range(len(shapes_sorted) - 1):
            a = shapes_sorted[i]
            b = shapes_sorted[i + 1]
            gap = b.left - (a.left + a.width)
            if 0 < gap < slide_w * 0.3:  # reasonable gap
                gutters.append(gap)

    avg_gutter = int(sum(gutters) / len(gutters)) if gutters else 0

    # Zone detection
    zones = []
    header_bottom = int(slide_h * 0.15)
    footer_top = int(slide_h * 0.85)

    zones.append({
        "name": "header",
        "left": dim(0), "top": dim(0),
        "width": dim(slide_w), "height": dim(header_bottom),
    })
    zones.append({
        "name": "content",
        "left": dim(min(lefts) if lefts else 0),
        "top": dim(header_bottom),
        "width": dim(slide_w - (min(lefts) if lefts else 0) - (min(rights) if rights else 0)),
        "height": dim(footer_top - header_bottom),
    })
    zones.append({
        "name": "footer",
        "left": dim(0), "top": dim(footer_top),
        "width": dim(slide_w), "height": dim(slide_h - footer_top),
    })

    return content_margins, {"columns": 0, "gutter": dim(avg_gutter), "detected_pattern": "auto-detected"}, zones


# ---------------------------------------------------------------------------
# Phase 4: Masters & Slide Layouts
# ---------------------------------------------------------------------------
def extract_slide_layouts(prs, major_font, minor_font):
    """Inventory all masters and slide layouts with placeholders and decorative shapes."""
    layouts_list = []

    for master in prs.slide_masters:
        master_name = master.name if hasattr(master, 'name') else "Default Master"
        for layout in master.slide_layouts:
            placeholders = []
            for ph in layout.placeholders:
                ph_type = "UNKNOWN"
                try:
                    ph_type = str(ph.placeholder_format.type).split('(')[0]
                except Exception:
                    pass
                placeholders.append({
                    "type": ph_type,
                    "idx": ph.placeholder_format.idx,
                    "left": dim(ph.left),
                    "top": dim(ph.top),
                    "width": dim(ph.width),
                    "height": dim(ph.height),
                })

            decorative = []
            for shape in layout.shapes:
                if shape.is_placeholder:
                    continue
                dec = {
                    "type": str(shape.shape_type).split('(')[0] if shape.shape_type else "unknown",
                    "fill": safe_fill_color(shape),
                    "position": {
                        "left_in": emu_to_inches(shape.left),
                        "top_in": emu_to_inches(shape.top),
                        "width_in": emu_to_inches(shape.width),
                        "height_in": emu_to_inches(shape.height),
                    },
                    "notes": "",
                }
                decorative.append(dec)

            # Background
            bg = None
            try:
                bg_el = layout.element.find(qn('p:cSld') + '/' + qn('p:bg'))
                if bg_el is not None:
                    srgb = bg_el.find('.//' + qn('a:srgbClr'))
                    if srgb is not None:
                        bg = f"#{srgb.get('val', '').upper()}"
                    else:
                        bg = "inherited"
                else:
                    bg = "inherited"
            except Exception:
                bg = "inherited"

            layouts_list.append({
                "name": layout.name,
                "master_name": master_name,
                "background": bg,
                "placeholders": placeholders,
                "decorative_shapes": decorative,
            })

    return layouts_list


# ---------------------------------------------------------------------------
# Phase 5: Asset Extraction
# ---------------------------------------------------------------------------
def extract_assets(pptx_path, output_dir):
    """Extract all media from the PPTX zip and classify them."""
    assets_dir = os.path.join(output_dir, "assets")
    os.makedirs(assets_dir, exist_ok=True)

    assets = []
    try:
        with zipfile.ZipFile(pptx_path, 'r') as z:
            media_files = [f for f in z.namelist() if f.startswith('ppt/media/')]
            for mf in media_files:
                filename = os.path.basename(mf)
                ext = os.path.splitext(filename)[1].lower()
                fmt_map = {'.png': 'png', '.jpg': 'jpg', '.jpeg': 'jpg', '.gif': 'gif',
                           '.svg': 'svg', '.emf': 'emf', '.wmf': 'wmf', '.tiff': 'tiff', '.bmp': 'bmp'}
                fmt = fmt_map.get(ext, ext.lstrip('.'))

                data = z.read(mf)
                file_size = len(data)

                # Save to assets dir
                out_path = os.path.join(assets_dir, filename)
                with open(out_path, 'wb') as f:
                    f.write(data)

                # Get dimensions
                w_px, h_px = 0, 0
                if HAS_PIL and fmt in ('png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff'):
                    try:
                        from io import BytesIO
                        img = Image.open(BytesIO(data))
                        w_px, h_px = img.size
                    except Exception:
                        pass

                # Classification heuristic
                classification = "photo"
                if w_px > 0 and h_px > 0:
                    if w_px < 200 and h_px < 200:
                        classification = "icon"
                    elif w_px < 400 and h_px < 200:
                        classification = "logo"
                    elif w_px > 900 and h_px > 500:
                        classification = "background"
                if fmt in ('emf', 'wmf', 'svg'):
                    classification = "vector_graphic"

                is_brand = classification in ('logo', 'icon', 'vector_graphic')

                assets.append({
                    "filename": filename,
                    "format": fmt,
                    "dimensions_px": {"width": w_px, "height": h_px},
                    "file_size_bytes": file_size,
                    "classification": classification,
                    "found_on_slides": [],  # Would need rels parsing for precision
                    "found_on_master_or_layout": "",
                    "extracted_to": os.path.abspath(out_path),
                    "is_brand_asset": is_brand,
                })
    except Exception as e:
        print(f"Warning: asset extraction error: {e}")

    return assets


# ---------------------------------------------------------------------------
# Phase 6: Components
# ---------------------------------------------------------------------------
def extract_components(prs):
    """Detect recurring components: dividers, bullets, tables, CTAs."""
    dividers = []
    bullets_set = set()
    table_styles = []
    chart_colors = []

    for slide in prs.slides:
        for shape in slide.shapes:
            # Divider detection: very thin shapes
            if shape.width and shape.height:
                w_pt = emu_to_pt(shape.width)
                h_pt = emu_to_pt(shape.height)
                if h_pt is not None and h_pt < 5 and w_pt and w_pt > 100:
                    dividers.append({
                        "orientation": "horizontal",
                        "color": safe_fill_color(shape) or safe_line_color(shape),
                        "thickness_pt": h_pt,
                        "position_pattern": f"slide at y={emu_to_inches(shape.top)}in",
                    })
                elif w_pt is not None and w_pt < 5 and h_pt and h_pt > 100:
                    dividers.append({
                        "orientation": "vertical",
                        "color": safe_fill_color(shape) or safe_line_color(shape),
                        "thickness_pt": w_pt,
                        "position_pattern": f"slide at x={emu_to_inches(shape.left)}in",
                    })

            # Bullet detection from text frames
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    try:
                        pPr = para._p.find(qn('a:pPr'))
                        if pPr is not None:
                            buChar = pPr.find(qn('a:buChar'))
                            if buChar is not None:
                                char = buChar.get('char', '•')
                                bullets_set.add(char)
                    except Exception:
                        pass

            # Table style extraction
            if shape.has_table:
                tbl = shape.table
                style_info = {"header_fill": None, "border_color": None, "row_count": len(tbl.rows)}
                # Header row fill
                if len(tbl.rows) > 0:
                    try:
                        cell = tbl.rows[0].cells[0]
                        tc_el = cell._tc
                        srgb = tc_el.find('.//' + qn('a:srgbClr'))
                        if srgb is not None:
                            style_info["header_fill"] = f"#{srgb.get('val', '').upper()}"
                    except Exception:
                        pass
                table_styles.append(style_info)

    # Deduplicate dividers
    seen = set()
    unique_dividers = []
    for d in dividers:
        key = (d["orientation"], d["color"], d["thickness_pt"])
        if key not in seen:
            seen.add(key)
            unique_dividers.append(d)

    bullets = [{"character": c, "color": None, "size_pt": None, "indent_inches": None} for c in bullets_set]

    return {
        "dividers": unique_dividers[:10],
        "bullets": bullets,
        "table_styles": table_styles[:5],
        "buttons_ctas": [],
        "cards": [],
        "chart_color_sequence": chart_colors,
    }


# ---------------------------------------------------------------------------
# Phase 7: Metadata & Slide Overrides
# ---------------------------------------------------------------------------
def extract_metadata(prs):
    """Extract core and custom properties."""
    props = prs.core_properties
    return {
        "author": props.author or "",
        "company": "",  # Not directly in core_properties
        "title": props.title or "",
        "subject": props.subject or "",
        "keywords": props.keywords or "",
        "created": str(props.created) if props.created else "",
        "modified": str(props.modified) if props.modified else "",
        "category": props.category or "",
    }


def extract_slide_overrides(prs):
    """Detect slides with custom backgrounds or unique layouts."""
    overrides = []
    for idx, slide in enumerate(prs.slides, 1):
        has_custom_bg = False
        bg_detail = "default"
        try:
            bg_el = slide.element.find(qn('p:cSld') + '/' + qn('p:bg'))
            if bg_el is not None:
                has_custom_bg = True
                srgb = bg_el.find('.//' + qn('a:srgbClr'))
                if srgb is not None:
                    bg_detail = f"#{srgb.get('val', '').upper()}"
                grad = bg_el.find('.//' + qn('a:gradFill'))
                if grad is not None:
                    bg_detail = "gradient"
                blip = bg_el.find('.//' + qn('a:blip'))
                if blip is not None:
                    bg_detail = "image"
        except Exception:
            pass

        if has_custom_bg:
            overrides.append({
                "slide_number": idx,
                "has_custom_background": True,
                "background_detail": bg_detail,
                "notes": "",
            })
    return overrides


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def extract_branding(pptx_path, output_dir="/mnt/user-data/outputs"):
    """Main extraction pipeline."""
    pptx_path = os.path.abspath(pptx_path)
    if not os.path.isfile(pptx_path):
        print(f"ERROR: File not found: {pptx_path}")
        sys.exit(1)

    print(f"Analyzing: {pptx_path}")
    prs = Presentation(pptx_path)
    os.makedirs(output_dir, exist_ok=True)

    # Phase 1: Colors
    print("Phase 1/7: Extracting theme colors...")
    theme_colors = extract_theme_colors(prs)
    major_font, minor_font = extract_theme_fonts(prs)
    fill_c, font_c, border_c = collect_all_colors(prs)
    dominant = build_dominant_colors(fill_c, font_c, border_c)
    gradients = extract_gradients(prs)

    # Phase 2: Typography
    print("Phase 2/7: Extracting typography...")
    text_styles = extract_typography(prs, major_font, minor_font)

    # Phase 3: Layout
    print("Phase 3/7: Analyzing layout & spatial...")
    content_margins, grid, zones = extract_layouts(prs)

    # Phase 4: Masters & Layouts
    print("Phase 4/7: Inventorying masters & layouts...")
    slide_layouts = extract_slide_layouts(prs, major_font, minor_font)

    # Phase 5: Assets
    print("Phase 5/7: Extracting assets...")
    assets = extract_assets(pptx_path, output_dir)

    # Phase 6: Components
    print("Phase 6/7: Detecting components...")
    components = extract_components(prs)

    # Phase 7: Metadata
    print("Phase 7/7: Reading metadata...")
    meta_props = extract_metadata(prs)
    slide_overrides = extract_slide_overrides(prs)

    # Assemble result
    result = {
        "meta": {
            "source_filename": os.path.basename(pptx_path),
            "slide_count": len(prs.slides),
            "slide_width": dim(prs.slide_width),
            "slide_height": dim(prs.slide_height),
            **meta_props,
        },
        "color_palette": {
            "theme_colors": theme_colors,
            "dominant_colors": dominant,
            "gradients": gradients,
        },
        "typography": {
            "theme_fonts": {
                "major": major_font,
                "minor": minor_font,
            },
            "text_styles": text_styles,
        },
        "layouts": {
            "content_margins": content_margins,
            "grid": grid,
            "zones": zones,
            "slide_layouts": slide_layouts,
        },
        "components": components,
        "assets": assets,
        "slide_overrides": slide_overrides,
    }

    # Save
    out_path = os.path.join(output_dir, "branding_template.json")
    with open(out_path, 'w') as f:
        json.dump(result, f, indent=2, ensure_ascii=False, default=str)

    print(f"\n✅ Done! Output saved to: {out_path}")
    print(f"   Assets saved to: {os.path.join(output_dir, 'assets')}/")
    asset_count = len(assets)
    brand_count = sum(1 for a in assets if a["is_brand_asset"])
    print(f"   {asset_count} images extracted ({brand_count} classified as brand assets)")
    print(f"   {len(text_styles)} text styles detected")
    print(f"   {len(dominant)} unique colors found")

    return result


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_branding.py <path_to.pptx> [--output-dir <dir>]")
        sys.exit(1)

    pptx_file = sys.argv[1]
    out_dir = "/mnt/user-data/outputs"
    if "--output-dir" in sys.argv:
        idx = sys.argv.index("--output-dir")
        if idx + 1 < len(sys.argv):
            out_dir = sys.argv[idx + 1]

    extract_branding(pptx_file, out_dir)
